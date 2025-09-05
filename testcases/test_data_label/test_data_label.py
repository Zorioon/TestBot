import asyncio
import json
import os
import time
from typing import Any, Dict, List, Optional
import pytest
import xlwt
from entity.api_asset.api_asset import ApiAssetLabelDetail, ApiAssetRecord
from entity.file_asset.file_asset import FileAssetRecord
from utils.file_tools.file_utils import FileUtils
from utils.file_tools.word_doc_utils import WordDocManager
from utils.file_tools.zip_utils import ZipUtils
from utils.log_tools.logger_utils import get_logger
from utils.request_tools.async_http_client import AsyncHttpClient
from utils.sr_tools.apione_utils import ApioneUtils
from utils.ssh_tools.ssh_connect import AsyncSSHClient
from utils.ssh_tools.ssh_operation import SSHOperation
from pathlib import Path
import pandas as pd
from docx import Document
from fpdf import FPDF
from pptx import Presentation
from win32com.client import constants
from win32com import client as win32
import pywintypes
from openpyxl import Workbook, load_workbook
from openpyxl.styles import Border, Side, PatternFill, Font, Alignment, colors

log = get_logger(__name__)


@pytest.fixture(scope="class")
def all_data_label_refers():
    file_path = FileUtils.find_file_from_root(
        "data/data_label/specification_refer.json"
    )
    with open(file_path, "r", encoding="utf-8") as f:
        data_label_refers = json.load(f)
    return data_label_refers


@pytest.fixture(scope="class")
def all_data_labels():
    file_path = FileUtils.find_file_from_root("data/data_label/base_data_label.json")
    with open(file_path, "r", encoding="utf-8") as f:
        all_data_labels = json.load(f)
    return all_data_labels


class TestDataLabel:
    """数据标签测试"""

    def setup_class(cls):
        cls.test_results = []
        cls.word_doc_manager = WordDocManager()

    def teardown_class(cls):
        cls.word_doc_manager.close()

    @pytest.fixture(scope="class")
    def apps(proxy_apps):
        pass

    @staticmethod
    def load_specification():
        file_path = FileUtils.find_file_from_root("data/data_label/specification.json")
        with open(file_path, "r", encoding="utf-8") as f:
            specification_list = json.load(f)
        return [(item["name"], item["id"]) for item in specification_list]

    async def choose_specification(
        self,
        https_req: AsyncHttpClient,
        sc_ssh_client: AsyncSSHClient,
        specification_name: str,
        specification_id: int,
    ):
        try:
            # 1. 清理系统
            await SSHOperation.exec_single_command(
                sc_ssh_client, "cd /opt/apione && ./bin/apione --clean"
            )
            log.success("清理脚本执行成功")

            # 2. 检查 apione 启动
            await SSHOperation.check_process_log(
                ssh_client=sc_ssh_client,
                process_name="apione",
                keyword='"http server listening at" address=127.0.0.1:29300',
            )
            log.success("apione 进程已成功启动")

            # 3. 初始化规则
            await ApioneUtils.initial_rule(https_req, specification_id)
            log.success(
                f"规则初始化完成 (specification_name={specification_name}, spec_id={specification_id})"
            )

            # 4、检查 ata 启动
            await SSHOperation.check_process_log(
                ssh_client=sc_ssh_client,
                process_name="ata",
                keyword="connect to nsqd server: 127.0.0.1:7150 completed",
            )
            log.success("ata 进程已成功启动")
            # await asyncio.sleep(30)
        except Exception:

            log.exception("初始化规则失败")
            raise

    def compare_api_asset_label_result(
        self,
        api_asset_data_label: Dict[str, Any],
        api_asset_label_detail: Optional[ApiAssetLabelDetail],
    ):
        """比较测试样本和实际匹配的情况，输出一个Excel表格

        Args:
            api_asset_data_label (Dict[str, Any]): _description_
            api_asset_label_detail (Optional[ApiAssetLabelDetail]): _description_
        """
        data_label_id = api_asset_data_label.get("id")
        data_label_name = api_asset_data_label.get("name")
        expected_start_line_data_label_value = (
            api_asset_data_label.get("start_line", None) or []
        )
        expected_headers_data_label_value = (
            api_asset_data_label.get("headers", None) or []
        )
        expected_body_data_label_value = api_asset_data_label.get("body", None) or []
        expected_count = (
            len(expected_start_line_data_label_value)
            + len(expected_headers_data_label_value)
            + len(expected_body_data_label_value)
        )
        results: Dict[str, Any] = {
            "request": {
                "id": data_label_id,
                "name": data_label_name,
                "sample": {
                    "start_line": expected_start_line_data_label_value,
                    "headers": expected_headers_data_label_value,
                    "body": expected_body_data_label_value,
                    "count": expected_count,
                },
                "matched": {
                    "start_line": None,
                    "headers": None,
                    "body": None,
                    "count": None,
                },
                "unmatched": {
                    "start_line": None,
                    "headers": None,
                    "body": None,
                    "count": None,
                },
                "misidentification": {
                    "start_line": None,
                    "headers": None,
                    "body": None,
                    "count": None,
                },
                "status": None,
            },
            "response": {
                "id": data_label_id,
                "name": data_label_name,
                "sample": {
                    "start_line": expected_start_line_data_label_value,
                    "headers": expected_headers_data_label_value,
                    "body": expected_body_data_label_value,
                    "count": expected_count,
                },
                "matched": {
                    "start_line": None,
                    "headers": None,
                    "body": None,
                    "count": None,
                },
                "unmatched": {
                    "start_line": None,
                    "headers": None,
                    "body": None,
                    "count": None,
                },
                "misidentification": {
                    "start_line": None,
                    "headers": None,
                    "body": None,
                    "count": None,
                },
                "status": None,
            },
        }
        for part in ["request", "response"]:
            real_data_label_value = getattr(api_asset_label_detail, part, None) or {}
            specified_data_label_in_start_line = (
                real_data_label_value.get("start_line", None) or {}
            ).get(data_label_name, {})
            specified_data_label_in_headers = (
                real_data_label_value.get("headers", None) or {}
            ).get(data_label_name, {})
            specified_data_label_in_body = (
                real_data_label_value.get("body", None) or {}
            ).get(data_label_name, {})
            matched_count = 0
            unmatched_count = 0
            misidentification_count = 0

            if specified_data_label_in_start_line:
                results[part]["matched"]["start_line"] = (
                    specified_data_label_in_start_line["contents"]
                )
                matched_count += specified_data_label_in_start_line["count"]
                results[part]["matched"]["count"] = matched_count
                unmatched_data_label_in_start_line = [
                    content
                    for content in expected_start_line_data_label_value
                    if content not in specified_data_label_in_start_line["contents"]
                ] or expected_start_line_data_label_value

                unmatched_count += len(unmatched_data_label_in_start_line)
                results[part]["unmatched"][
                    "start_line"
                ] = unmatched_data_label_in_start_line
                results[part]["unmatched"]["count"] = unmatched_count

                del real_data_label_value.get("start_line")[data_label_name]
                results[part]["misidentification"]["start_line"] = (
                    real_data_label_value.get("start_line")
                )
                misidentification_count += len(real_data_label_value.get("start_line"))
                results[part]["misidentification"]["count"] = misidentification_count

            if specified_data_label_in_headers:
                results[part]["matched"]["headers"] = specified_data_label_in_headers[
                    "contents"
                ]
                matched_count += specified_data_label_in_headers["count"]
                results[part]["matched"]["count"] = matched_count
                unmatched_data_label_in_headers = [
                    content
                    for content in expected_headers_data_label_value
                    if content not in specified_data_label_in_headers["contents"]
                ] or expected_headers_data_label_value

                unmatched_count += len(unmatched_data_label_in_headers)
                results[part]["unmatched"]["headers"] = unmatched_data_label_in_headers
                results[part]["unmatched"]["count"] = unmatched_count

                del real_data_label_value.get("headers")[data_label_name]
                results[part]["misidentification"]["headers"] = (
                    real_data_label_value.get("headers")
                )
                misidentification_count += len(real_data_label_value.get("headers"))
                results[part]["misidentification"]["count"] = misidentification_count

            if specified_data_label_in_body:
                results[part]["matched"]["body"] = specified_data_label_in_body[
                    "contents"
                ]
                matched_count += specified_data_label_in_body["count"]
                results[part]["matched"]["count"] = matched_count
                unmatched_data_label_in_body = [
                    content
                    for content in expected_body_data_label_value
                    if content not in specified_data_label_in_body["contents"]
                ] or expected_body_data_label_value
                unmatched_count += len(unmatched_data_label_in_body)
                results[part]["unmatched"]["body"] = unmatched_data_label_in_body
                results[part]["unmatched"]["count"] = unmatched_count

                del real_data_label_value.get("body")[data_label_name]
                results[part]["misidentification"]["body"] = real_data_label_value.get(
                    "body"
                )
                misidentification_count += len(real_data_label_value.get("body"))
                results[part]["misidentification"]["count"] = misidentification_count

            results[part]["status"] = (
                "PASS"
                if matched_count == expected_count
                and unmatched_count == 0
                and misidentification_count == 0
                else "FAILED"
            )
        return results

    async def compare_file_asset_label_result(
        self,
        https_req: AsyncHttpClient,
        file_asset_data_label: Dict[str, Any],
        file_lists: List[str],
    ):
        data_label_id = file_asset_data_label.get("id")
        data_label_name = file_asset_data_label.get("name")
        expected_file_data_label_value = (
            file_asset_data_label.get("file_data", None) or []
        )
        all_passed = True
        expected_count = len(expected_file_data_label_value)
        file_data_label_result = {
            "id": data_label_id,
            "name": data_label_name,
            # ".doc": None,
            ".docx": None,
            ".xls": None,
            ".xlsx": None,
            ".txt": None,
            ".pptx": None,
            ".pdf": None,
            ".csv": None,
            # ".zip": None,
            "status": None,
        }

        for file_path in file_lists:
            file_md5 = FileUtils.calculate_file_md5(file_path)
            file_name = os.path.basename(file_path)
            file_ext = os.path.splitext(file_path)[1]
            file_asset: Optional[FileAssetRecord] = (
                await ApioneUtils.get_file_asset_record(https_req, file_name, file_md5)
            )
            file_data_label_detail = await ApioneUtils.get_file_asset_label_detail(
                https_req, file_asset.id
            )
            matched_count = file_data_label_detail.get(data_label_name, 0)
            misidentification = 0
            if data_label_name in file_data_label_detail:
                del file_data_label_detail[data_label_name]
                misidentification = file_data_label_detail
            else:
                misidentification = file_data_label_detail
            file_data_label_result[file_ext] = {
                "target_file": file_name,
                "expected_count": expected_count,
                "matched_count": matched_count,
                "misidentification": misidentification,
            }
            if expected_count != matched_count or len(misidentification):
                all_passed = False
        file_data_label_result["status"] = "PASS" if all_passed else "FAILED"

        return file_data_label_result

    # def export_api_label_to_excel(
    #     self,
    #     data_label_test_result: List[Dict[str, Dict]],
    #     specification_name: str,
    # ):
    #     """导出数据到Excel，分别创建request和response的sheet"""
    #     file_path = FileUtils.find_file_from_root(f"files/api_label_test_result_{specification_name}.xlsx", create_if_not_exists=True)

    #     def handle_empty(x):
    #         # 处理None、空值、0
    #         if x in (None, 0, "") or (hasattr(x, "__len__") and len(x) == 0):
    #             return "-"
    #         # 处理列表和字典，转换为JSON字符串
    #         elif isinstance(x, (list, dict)):
    #             try:
    #                 json_str = json.dumps(x, ensure_ascii=False, indent=2)
    #                 if len(json_str) > 1000:  # 限制最大长度
    #                     return json.dumps(x, ensure_ascii=False)[:1000] + "..."
    #                 return json_str
    #             except:
    #                 return str(x)
    #         else:
    #             return x

    #     # ===== 样式定义 =====
    #     # 边框样式
    #     thin_border = Border(
    #         left=Side(style="thin"),
    #         right=Side(style="thin"),
    #         top=Side(style="thin"),
    #         bottom=Side(style="thin"),
    #     )

    #     # 表头
    #     header_fill = PatternFill("solid", fgColor="FF4F81BD")  # 蓝色底（更醒目）
    #     header_font = Font(bold=True, color="FFFFFFFF", size=12)  # 白字，加大字号

    #     # 状态列
    #     pass_fill = PatternFill("solid", fgColor="92D050")  # 绿色底
    #     pass_font = Font(color="006100", bold=True)  # 深绿字
    #     fail_fill = PatternFill("solid", fgColor="FF0000")  # 红色底
    #     fail_font = Font(color="FFFFFF", bold=True)  # 白字

    #     # count 列
    #     count_blue_font = Font(color="0000FF", bold=True)  # 蓝字
    #     count_gray_font = Font(color="808080")  # 灰字

    #     # 类型列底色（使用更明显的颜色区分）
    #     type_fills = {
    #         "样本": PatternFill("solid", fgColor="FFFFCC"),  # 浅黄色
    #         "已匹配": PatternFill("solid", fgColor="E2EFDA"),  # 浅绿色
    #         "未匹配": PatternFill("solid", fgColor="FCE4D6"),  # 浅橙色
    #         "误匹配": PatternFill("solid", fgColor="FFE6E6"),  # 浅红色
    #     }

    #     # 交替行颜色（增强可读性）
    #     even_row_fill = PatternFill("solid", fgColor="F2F2F2")  # 浅灰色
    #     odd_row_fill = PatternFill("solid", fgColor="FFFFFF")  # 白色

    #     try:
    #         wb = load_workbook(file_path)
    #     except FileNotFoundError:
    #         wb = Workbook()
    #         if "Sheet" in wb.sheetnames:
    #             wb.remove(wb["Sheet"])

    #     # 处理每个sheet类型
    #     for sheet_type in ["request", "response"]:
    #         sheet_name = (
    #             "请求详情识别结果" if sheet_type == "request" else "响应详情识别结果"
    #         )

    #         if sheet_name in wb.sheetnames:
    #             ws = wb[sheet_name]
    #             # 删除数据行但保留表头
    #             if ws.max_row > 1:
    #                 ws.delete_rows(2, ws.max_row - 1)
    #             # 确保表头样式存在
    #             if ws.max_row >= 1:
    #                 for cell in ws[1]:
    #                     cell.fill = header_fill
    #                     cell.font = header_font
    #                     cell.alignment = Alignment(
    #                         horizontal="center", vertical="center"
    #                     )
    #                     cell.border = thin_border
    #         else:
    #             ws = wb.create_sheet(sheet_name)
    #             ws.append(
    #                 [
    #                     "编号",
    #                     "名称",
    #                     "类型",
    #                     "start_line",
    #                     "headers",
    #                     "body",
    #                     "匹配数量",
    #                     "状态",
    #                 ]
    #             )
    #             # 设置表头样式
    #             for cell in ws[1]:
    #                 cell.fill = header_fill
    #                 cell.font = header_font
    #                 cell.alignment = Alignment(horizontal="center", vertical="center")
    #                 cell.border = thin_border

    #         # 处理该类型的所有数据项
    #         for item_index, item in enumerate(data_label_test_result):
    #             if sheet_type in item:
    #                 sheet_data = item[sheet_type]

    #                 types = ["样本", "已匹配", "未匹配", "误匹配"]
    #                 key_map = {
    #                     "样本": "sample",
    #                     "已匹配": "matched",
    #                     "未匹配": "unmatched",
    #                     "误匹配": "misidentification",
    #                 }

    #                 start_row = ws.max_row + 1

    #                 for t in types:
    #                     key = key_map[t]
    #                     row_data = sheet_data.get(key, {})

    #                     row = [
    #                         handle_empty(sheet_data.get("id")),
    #                         handle_empty(sheet_data.get("name")),
    #                         t,
    #                         handle_empty(row_data.get("start_line")),
    #                         handle_empty(row_data.get("headers")),
    #                         handle_empty(row_data.get("body")),
    #                         handle_empty(row_data.get("count", 0)),
    #                         handle_empty(sheet_data.get("status")),
    #                     ]
    #                     ws.append(row)

    #                     # 获取当前行
    #                     cur_row = ws.max_row

    #                     # === 设置整行边框 ===
    #                     for col in range(1, 9):  # 1到8列
    #                         ws.cell(cur_row, col).border = thin_border

    #                     # === 交替行颜色 ===
    #                     row_fill = even_row_fill if cur_row % 2 == 0 else odd_row_fill
    #                     for col in range(1, 9):
    #                         if col not in [3, 7, 8]:  # 排除已经有特殊样式的列
    #                             ws.cell(cur_row, col).fill = row_fill

    #                     # === 类型列底色 ===
    #                     type_cell = ws.cell(cur_row, 3)
    #                     if t in type_fills:
    #                         type_cell.fill = type_fills[t]
    #                     type_cell.border = thin_border
    #                     type_cell.alignment = Alignment(
    #                         horizontal="center", vertical="center"
    #                     )

    #                     # === 状态列样式 ===
    #                     status_cell = ws.cell(cur_row, 8)
    #                     if status_cell.value == "PASS":
    #                         status_cell.fill = pass_fill
    #                         status_cell.font = pass_font
    #                     elif status_cell.value == "FAILED":
    #                         status_cell.fill = fail_fill
    #                         status_cell.font = fail_font
    #                     status_cell.border = thin_border
    #                     status_cell.alignment = Alignment(
    #                         horizontal="center", vertical="center"
    #                     )

    #                     # === count 列样式 ===
    #                     count_cell = ws.cell(cur_row, 7)
    #                     try:
    #                         if count_cell.value not in ("-", None, ""):
    #                             if int(count_cell.value) > 0:
    #                                 count_cell.font = count_blue_font
    #                             else:
    #                                 count_cell.font = count_gray_font
    #                         else:
    #                             count_cell.font = count_gray_font
    #                     except:
    #                         count_cell.font = count_gray_font
    #                     count_cell.border = thin_border
    #                     count_cell.alignment = Alignment(
    #                         horizontal="center", vertical="center"
    #                     )

    #                     # === 编号和名称列居中 ===
    #                     for col in [1, 2]:
    #                         ws.cell(cur_row, col).alignment = Alignment(
    #                             horizontal="center", vertical="center"
    #                         )

    #                 end_row = ws.max_row

    #                 # 合并单元格（编号、名称、status）
    #                 for col in [1, 2, 8]:
    #                     ws.merge_cells(
    #                         start_row=start_row,
    #                         start_column=col,
    #                         end_row=end_row,
    #                         end_column=col,
    #                     )
    #                     # 设置合并后单元格的样式
    #                     merged_cell = ws.cell(start_row, col)
    #                     merged_cell.alignment = Alignment(
    #                         vertical="center", horizontal="center"
    #                     )
    #                     merged_cell.border = thin_border

    #                     # 为合并的单元格设置交替行颜色
    #                     if col not in [3, 7, 8]:
    #                         row_fill = (
    #                             even_row_fill if start_row % 2 == 0 else odd_row_fill
    #                         )
    #                         merged_cell.fill = row_fill

    #     # 调整列宽
    #     for sheet_name in wb.sheetnames:
    #         if sheet_name in ["请求详情识别结果", "响应详情识别结果"]:
    #             ws = wb[sheet_name]
    #             # 设置固定列宽，避免自动调整过于紧凑
    #             column_widths = {
    #                 "A": 15,  # 编号
    #                 "B": 20,  # 名称
    #                 "C": 10,  # 类型
    #                 "D": 12,  # start_line
    #                 "E": 25,  # headers
    #                 "F": 25,  # body
    #                 "G": 10,  # count
    #                 "H": 12,  # status
    #             }

    #             for col_letter, width in column_widths.items():
    #                 ws.column_dimensions[col_letter].width = width

    #             # 设置行高
    #             for row in range(2, ws.max_row + 1):
    #                 ws.row_dimensions[row].height = 25

    #     wb.save(file_path)
    #     log.success(f"数据已导出到 {file_path}")

    # def export_file_label_to_excel(
    #     self,
    #     file_asset_data_label_test_result: List[Dict],
    #     specification_name: str
    # ):
    #     """导出文件标签测试结果到Excel"""
    #     file_path = FileUtils.find_file_from_root(f"files/file_label_test_result/file_label_test_result_{specification_name}.xlsx", create_if_not_exists=True)

    #     def handle_empty(x):
    #         if x in (None, 0, "") or (hasattr(x, "__len__") and len(x) == 0):
    #             return "-"
    #         elif isinstance(x, (list, dict)):
    #             try:
    #                 json_str = json.dumps(x, ensure_ascii=False, indent=2)
    #                 if len(json_str) > 1000:
    #                     return json.dumps(x, ensure_ascii=False)[:1000] + "..."
    #                 return json_str
    #             except:
    #                 return str(x)
    #         else:
    #             return x

    #     # ===== 样式定义 =====
    #     thin_border = Border(
    #         left=Side(style="thin"),
    #         right=Side(style="thin"),
    #         top=Side(style="thin"),
    #         bottom=Side(style="thin"),
    #     )

    #     header_fill = PatternFill("solid", fgColor="FF4F81BD")
    #     header_font = Font(bold=True, color="FFFFFFFF", size=12)

    #     pass_fill = PatternFill("solid", fgColor="92D050")
    #     pass_font = Font(color="006100", bold=True)
    #     fail_fill = PatternFill("solid", fgColor="FF0000")
    #     fail_font = Font(color="FFFFFF", bold=True)

    #     # count 列样式
    #     count_blue_font = Font(color="0000FF", bold=True)  # 蓝字
    #     count_gray_font = Font(color="808080")  # 灰字

    #     file_type_fills = {
    #         ".doc": PatternFill("solid", fgColor="E2EFDA"),
    #         ".docx": PatternFill("solid", fgColor="DDEBF7"),
    #         ".xls": PatternFill("solid", fgColor="FFF2CC"),
    #         ".xlsx": PatternFill("solid", fgColor="FCE4D6"),
    #         ".txt": PatternFill("solid", fgColor="E2EFDA"),
    #         ".pptx": PatternFill("solid", fgColor="DDEBF7"),
    #         ".pdf": PatternFill("solid", fgColor="FFF2CC"),
    #         ".csv": PatternFill("solid", fgColor="FCE4D6"),
    #         # ".zip": PatternFill("solid", fgColor="FFE6E6"),
    #     }

    #     even_row_fill = PatternFill("solid", fgColor="F2F2F2")
    #     odd_row_fill = PatternFill("solid", fgColor="FFFFFF")

    #     try:
    #         wb = load_workbook(file_path)
    #     except FileNotFoundError:
    #         wb = Workbook()
    #         if "Sheet" in wb.sheetnames:
    #             wb.remove(wb["Sheet"])

    #     sheet_name = "文件标签识别结果"
    #     if sheet_name in wb.sheetnames:
    #         ws = wb[sheet_name]
    #         if ws.max_row > 1:
    #             ws.delete_rows(2, ws.max_row - 1)
    #         for cell in ws[1]:
    #             cell.fill = header_fill
    #             cell.font = header_font
    #             cell.alignment = Alignment(horizontal="center", vertical="center")
    #             cell.border = thin_border
    #     else:
    #         ws = wb.create_sheet(sheet_name)
    #         ws.append(
    #             [
    #                 "编号",
    #                 "名称",
    #                 "文件类型",
    #                 "目标文件",
    #                 "预期数量",
    #                 "匹配数量",
    #                 "误匹配详情",
    #                 "状态",
    #             ]
    #         )
    #         for cell in ws[1]:
    #             cell.fill = header_fill
    #             cell.font = header_font
    #             cell.alignment = Alignment(horizontal="center", vertical="center")
    #             cell.border = thin_border

    #     file_types = [
    #         ".doc",
    #         ".docx",
    #         ".xls",
    #         ".xlsx",
    #         ".txt",
    #         ".pptx",
    #         ".pdf",
    #         ".csv",
    #         # ".zip",
    #     ]

    #     for item_index, item in enumerate(file_asset_data_label_test_result):
    #         start_row = ws.max_row + 1
    #         for file_type in file_types:
    #             file_data = item.get(file_type)
    #             row = [
    #                 handle_empty(item.get("id")),
    #                 handle_empty(item.get("name")),
    #                 file_type,
    #                 handle_empty(file_data.get("target_file") if file_data else None),
    #                 handle_empty(
    #                     file_data.get("expected_count") if file_data else None
    #                 ),
    #                 handle_empty(file_data.get("matched_count") if file_data else None),
    #                 handle_empty(
    #                     file_data.get("misidentification") if file_data else None
    #                 ),
    #                 handle_empty(item.get("status")),
    #             ]
    #             ws.append(row)
    #             cur_row = ws.max_row

    #             # 设置整行边框和基础样式
    #             for col in range(1, 9):
    #                 cell = ws.cell(cur_row, col)
    #                 cell.border = thin_border
    #                 # 交替行颜色
    #                 if col not in [3, 6, 8]:  # 排除文件类型、匹配数量、状态列
    #                     cell.fill = even_row_fill if cur_row % 2 == 0 else odd_row_fill

    #             # 文件类型列
    #             type_cell = ws.cell(cur_row, 3)
    #             if file_type in file_type_fills:
    #                 type_cell.fill = file_type_fills[file_type]
    #             type_cell.alignment = Alignment(horizontal="center", vertical="center")

    #             # 匹配数量列样式
    #             count_cell = ws.cell(cur_row, 6)
    #             try:
    #                 if count_cell.value not in ("-", None, ""):
    #                     if int(count_cell.value) > 0:
    #                         count_cell.font = count_blue_font
    #                     else:
    #                         count_cell.font = count_gray_font
    #                 else:
    #                     count_cell.font = count_gray_font
    #             except:
    #                 count_cell.font = count_gray_font
    #             count_cell.alignment = Alignment(horizontal="right", vertical="center")

    #             # 状态列
    #             status_cell = ws.cell(cur_row, 8)
    #             status_value = str(status_cell.value or "").upper()
    #             if "PASS" in status_value:
    #                 status_cell.fill = pass_fill
    #                 status_cell.font = pass_font
    #             elif "FAILED" in status_value:
    #                 status_cell.fill = fail_fill
    #                 status_cell.font = fail_font
    #             status_cell.alignment = Alignment(
    #                 horizontal="center", vertical="center"
    #             )

    #             # 编号和名称居中
    #             for col in [1, 2]:
    #                 ws.cell(cur_row, col).alignment = Alignment(
    #                     horizontal="center", vertical="center"
    #                 )
    #             # 预期数量列右对齐
    #             ws.cell(cur_row, 5).alignment = Alignment(
    #                 horizontal="right", vertical="center"
    #             )

    #         end_row = ws.max_row

    #         # 合并单元格
    #         for col in [1, 2, 8]:
    #             ws.merge_cells(
    #                 start_row=start_row,
    #                 start_column=col,
    #                 end_row=end_row,
    #                 end_column=col,
    #             )
    #             merged_cell = ws.cell(start_row, col)
    #             merged_cell.alignment = Alignment(
    #                 horizontal="center", vertical="center"
    #             )
    #             merged_cell.border = thin_border
    #             # 设置背景色
    #             if col != 8:  # 编号和名称列使用交替行色
    #                 merged_cell.fill = (
    #                     even_row_fill if start_row % 2 == 0 else odd_row_fill
    #                 )
    #             else:  # 状态列
    #                 status_value = str(merged_cell.value or "").upper()
    #                 if "PASS" in status_value:
    #                     merged_cell.fill = pass_fill
    #                     merged_cell.font = pass_font
    #                 elif "FAILED" in status_value:
    #                     merged_cell.fill = fail_fill
    #                     merged_cell.font = fail_font

    #     # 调整列宽
    #     widths = {
    #         "A": 15,
    #         "B": 20,
    #         "C": 12,
    #         "D": 30,
    #         "E": 12,
    #         "F": 12,
    #         "G": 25,
    #         "H": 12,
    #     }
    #     for col, width in widths.items():
    #         ws.column_dimensions[col].width = width
    #     for row in range(2, ws.max_row + 1):
    #         ws.row_dimensions[row].height = 25

    #     wb.save(file_path)
    #     log.success(f"文件标签测试结果已导出到 {file_path}")

    # def export_label_to_excel(
    #     self,
    #     api_asset_data_label_test_result: List[Dict[str, Dict]],
    #     file_asset_data_label_test_result: List[Dict],
    #     specification_name: str
    # ):
    #     """导出 API 和文件标签测试结果到同一个 Excel 文件"""
    #     file_path = FileUtils.find_file_from_root(
    #         f"files/label_test_result_{specification_name}.xlsx",
    #         create_if_not_exists=True
    #     )

    #     def handle_empty(x):
    #         if x in (None, 0, "") or (hasattr(x, "__len__") and len(x) == 0):
    #             return "-"
    #         elif isinstance(x, (list, dict)):
    #             try:
    #                 json_str = json.dumps(x, ensure_ascii=False, indent=2)
    #                 return json_str if len(json_str) <= 1000 else json.dumps(x, ensure_ascii=False)[:1000] + "..."
    #             except:
    #                 return str(x)
    #         else:
    #             return x

    #     # ===== 样式定义 =====
    #     thin_border = Border(
    #         left=Side(style="thin"), right=Side(style="thin"),
    #         top=Side(style="thin"), bottom=Side(style="thin")
    #     )
    #     header_fill = PatternFill("solid", fgColor="FF4F81BD")
    #     header_font = Font(bold=True, color="FFFFFFFF", size=12)
    #     pass_fill = PatternFill("solid", fgColor="92D050")
    #     pass_font = Font(color="006100", bold=True)
    #     fail_fill = PatternFill("solid", fgColor="FF0000")
    #     fail_font = Font(color="FFFFFF", bold=True)
    #     count_blue_font = Font(color="0000FF", bold=True)
    #     count_gray_font = Font(color="808080")
    #     even_row_fill = PatternFill("solid", fgColor="F2F2F2")
    #     odd_row_fill = PatternFill("solid", fgColor="FFFFFF")

    #     # API 类型颜色
    #     api_type_fills = {
    #         "样本": PatternFill("solid", fgColor="FFFFCC"),
    #         "已匹配": PatternFill("solid", fgColor="E2EFDA"),
    #         "未匹配": PatternFill("solid", fgColor="FCE4D6"),
    #         "误匹配": PatternFill("solid", fgColor="FFE6E6"),
    #     }

    #     # 文件类型颜色
    #     file_type_fills = {
    #         ".doc": PatternFill("solid", fgColor="E2EFDA"),
    #         ".docx": PatternFill("solid", fgColor="DDEBF7"),
    #         ".xls": PatternFill("solid", fgColor="FFF2CC"),
    #         ".xlsx": PatternFill("solid", fgColor="FCE4D6"),
    #         ".txt": PatternFill("solid", fgColor="E2EFDA"),
    #         ".pptx": PatternFill("solid", fgColor="DDEBF7"),
    #         ".pdf": PatternFill("solid", fgColor="FFF2CC"),
    #         ".csv": PatternFill("solid", fgColor="FCE4D6"),
    #         # ".zip": PatternFill("solid", fgColor="FFE6E6"),
    #     }

    #     try:
    #         wb = load_workbook(file_path)
    #     except FileNotFoundError:
    #         wb = Workbook()
    #         if "Sheet" in wb.sheetnames:
    #             wb.remove(wb["Sheet"])

    #     # ===== API 数据导出 =====
    #     for sheet_type in ["request", "response"]:
    #         sheet_name = "请求详情识别结果" if sheet_type == "request" else "响应详情识别结果"
    #         ws = wb[sheet_name] if sheet_name in wb.sheetnames else wb.create_sheet(sheet_name)
    #         if ws.max_row <= 1:
    #             ws.append(["编号","名称","类型","start_line","headers","body","匹配数量","状态"])
    #         for cell in ws[1]:
    #             cell.fill = header_fill
    #             cell.font = header_font
    #             cell.alignment = Alignment(horizontal="center", vertical="center")
    #             cell.border = thin_border

    #         for item in api_asset_data_label_test_result:
    #             if sheet_type not in item:
    #                 continue
    #             sheet_data = item[sheet_type]
    #             types = ["样本","已匹配","未匹配","误匹配"]
    #             key_map = {"样本":"sample","已匹配":"matched","未匹配":"unmatched","误匹配":"misidentification"}
    #             start_row = ws.max_row + 1
    #             for t in types:
    #                 row_data = sheet_data.get(key_map[t], {})
    #                 row = [
    #                     handle_empty(sheet_data.get("id")),
    #                     handle_empty(sheet_data.get("name")),
    #                     t,
    #                     handle_empty(row_data.get("start_line")),
    #                     handle_empty(row_data.get("headers")),
    #                     handle_empty(row_data.get("body")),
    #                     handle_empty(row_data.get("count",0)),
    #                     handle_empty(sheet_data.get("status"))
    #                 ]
    #                 ws.append(row)
    #                 cur_row = ws.max_row
    #                 # 设置样式
    #                 for col in range(1,9):
    #                     cell = ws.cell(cur_row,col)
    #                     cell.border = thin_border
    #                     if col not in [3,7,8]:
    #                         cell.fill = even_row_fill if cur_row%2==0 else odd_row_fill
    #                 ws.cell(cur_row,3).fill = api_type_fills[t]
    #                 ws.cell(cur_row,3).alignment = Alignment(horizontal="center", vertical="center")
    #                 count_cell = ws.cell(cur_row,7)
    #                 try:
    #                     if count_cell.value not in ("-",None,""):
    #                         count_cell.font = count_blue_font if int(count_cell.value)>0 else count_gray_font
    #                     else:
    #                         count_cell.font = count_gray_font
    #                 except:
    #                     count_cell.font = count_gray_font
    #                 count_cell.alignment = Alignment(horizontal="center", vertical="center")
    #                 status_cell = ws.cell(cur_row,8)
    #                 if status_cell.value=="PASS":
    #                     status_cell.fill, status_cell.font = pass_fill, pass_font
    #                 elif status_cell.value=="FAILED":
    #                     status_cell.fill, status_cell.font = fail_fill, fail_font
    #                 status_cell.alignment = Alignment(horizontal="center", vertical="center")
    #             # 合并编号、名称、状态
    #             end_row = ws.max_row
    #             for col in [1,2,8]:
    #                 ws.merge_cells(start_row=start_row, start_column=col, end_row=end_row, end_column=col)
    #                 merged_cell = ws.cell(start_row,col)
    #                 merged_cell.alignment = Alignment(horizontal="center", vertical="center")
    #                 merged_cell.border = thin_border
    #                 if col !=8:
    #                     merged_cell.fill = even_row_fill if start_row%2==0 else odd_row_fill
    #                 else:
    #                     if str(merged_cell.value or "").upper()=="PASS":
    #                         merged_cell.fill, merged_cell.font = pass_fill, pass_font
    #                     elif str(merged_cell.value or "").upper()=="FAILED":
    #                         merged_cell.fill, merged_cell.font = fail_fill, fail_font

    #     # ===== 文件标签数据导出（第三个 sheet） =====
    #     sheet_name = "文件标签识别结果"
    #     ws = wb[sheet_name] if sheet_name in wb.sheetnames else wb.create_sheet(sheet_name)
    #     if ws.max_row <= 1:
    #         ws.append(["编号","名称","文件类型","目标文件","预期数量","匹配数量","误匹配详情","状态"])
    #     for cell in ws[1]:
    #         cell.fill = header_fill
    #         cell.font = header_font
    #         cell.alignment = Alignment(horizontal="center", vertical="center")
    #         cell.border = thin_border

    #     file_types = list(file_type_fills.keys())
    #     for item in file_asset_data_label_test_result:
    #         start_row = ws.max_row + 1
    #         for ft in file_types:
    #             file_data = item.get(ft)
    #             row = [
    #                 handle_empty(item.get("id")),
    #                 handle_empty(item.get("name")),
    #                 ft,
    #                 handle_empty(file_data.get("target_file") if file_data else None),
    #                 handle_empty(file_data.get("expected_count") if file_data else None),
    #                 handle_empty(file_data.get("matched_count") if file_data else None),
    #                 handle_empty(file_data.get("misidentification") if file_data else None),
    #                 handle_empty(item.get("status"))
    #             ]
    #             ws.append(row)
    #             cur_row = ws.max_row
    #             for col in range(1,9):
    #                 cell = ws.cell(cur_row,col)
    #                 cell.border = thin_border
    #                 if col not in [3,6,8]:
    #                     cell.fill = even_row_fill if cur_row%2==0 else odd_row_fill
    #             ws.cell(cur_row,3).fill = file_type_fills.get(ft, odd_row_fill)
    #             ws.cell(cur_row,3).alignment = Alignment(horizontal="center", vertical="center")
    #             # count 列
    #             count_cell = ws.cell(cur_row,6)
    #             try:
    #                 if count_cell.value not in ("-",None,""):
    #                     count_cell.font = count_blue_font if int(count_cell.value)>0 else count_gray_font
    #                 else:
    #                     count_cell.font = count_gray_font
    #             except:
    #                 count_cell.font = count_gray_font
    #             count_cell.alignment = Alignment(horizontal="right", vertical="center")
    #             # 状态列
    #             status_cell = ws.cell(cur_row,8)
    #             status_value = str(status_cell.value or "").upper()
    #             if "PASS" in status_value:
    #                 status_cell.fill, status_cell.font = pass_fill, pass_font
    #             elif "FAILED" in status_value:
    #                 status_cell.fill, status_cell.font = fail_fill, fail_font
    #             status_cell.alignment = Alignment(horizontal="center", vertical="center")
    #         end_row = ws.max_row
    #         for col in [1,2,8]:
    #             ws.merge_cells(start_row=start_row, start_column=col, end_row=end_row, end_column=col)
    #             merged_cell = ws.cell(start_row,col)
    #             merged_cell.alignment = Alignment(horizontal="center", vertical="center")
    #             merged_cell.border = thin_border
    #             if col !=8:
    #                 merged_cell.fill = even_row_fill if start_row%2==0 else odd_row_fill
    #             else:
    #                 status_value = str(merged_cell.value or "").upper()
    #                 if "PASS" in status_value:
    #                     merged_cell.fill, merged_cell.font = pass_fill, pass_font
    #                 elif "FAILED" in status_value:
    #                     merged_cell.fill, merged_cell.font = fail_fill, fail_font

    #     # ===== 调整列宽和行高 =====
    #     sheet_columns = {
    #         "请求详情识别结果":{"A":15,"B":20,"C":10,"D":12,"E":25,"F":25,"G":10,"H":12},
    #         "响应详情识别结果":{"A":15,"B":20,"C":10,"D":12,"E":25,"F":25,"G":10,"H":12},
    #         "文件标签识别结果":{"A":15,"B":20,"C":12,"D":30,"E":12,"F":12,"G":25,"H":12},
    #     }
    #     for sheet_name, widths in sheet_columns.items():
    #         ws = wb[sheet_name]
    #         for col, width in widths.items():
    #             ws.column_dimensions[col].width = width
    #         for row in range(2, ws.max_row+1):
    #             ws.row_dimensions[row].height = 25

    #     wb.save(file_path)
    #     log.success(f"所有测试结果已导出到 {file_path}")

    def export_label_to_excel(
        self,
        api_asset_data_label_test_result: List[Dict[str, Dict]],
        file_asset_data_label_test_result: List[Dict],
        specification_name: str,
    ):
        """导出 API 和文件标签测试结果到同一个 Excel 文件，采用一次性覆盖逻辑"""
        file_path = FileUtils.find_file_from_root(
            f"files/data_label_file/test_result/{specification_name.replace('/', '_')}.xlsx",
            create_if_not_exists=True,
        )

        def handle_empty(x):
            if x in (None, 0, "") or (hasattr(x, "__len__") and len(x) == 0):
                return "-"
            elif isinstance(x, (list, dict)):
                try:
                    json_str = json.dumps(x, ensure_ascii=False, indent=2)
                    return (
                        json_str
                        if len(json_str) <= 1000
                        else json.dumps(x, ensure_ascii=False)
                    )
                except:
                    return str(x)
            else:
                return x

        # ===== 样式定义 =====
        thin_border = Border(
            left=Side(style="thin"),
            right=Side(style="thin"),
            top=Side(style="thin"),
            bottom=Side(style="thin"),
        )
        header_fill = PatternFill("solid", fgColor="FF4F81BD")
        header_font = Font(bold=True, color="FFFFFFFF", size=12)
        pass_fill = PatternFill("solid", fgColor="92D050")
        pass_font = Font(color="006100", bold=True)
        fail_fill = PatternFill("solid", fgColor="FF0000")
        fail_font = Font(color="FFFFFF", bold=True)
        count_blue_font = Font(color="0000FF", bold=True)
        count_gray_font = Font(color="808080")
        even_row_fill = PatternFill("solid", fgColor="F2F2F2")
        odd_row_fill = PatternFill("solid", fgColor="FFFFFF")

        # API 类型颜色
        api_type_fills = {
            "样本": PatternFill("solid", fgColor="FFFFCC"),
            "已匹配": PatternFill("solid", fgColor="E2EFDA"),
            "未匹配": PatternFill("solid", fgColor="FCE4D6"),
            "误匹配": PatternFill("solid", fgColor="FFE6E6"),
        }

        # 文件类型颜色
        file_type_fills = {
            # ".doc": PatternFill("solid", fgColor="E2EFDA"),
            ".docx": PatternFill("solid", fgColor="DDEBF7"),
            ".xls": PatternFill("solid", fgColor="FFF2CC"),
            ".xlsx": PatternFill("solid", fgColor="FCE4D6"),
            ".txt": PatternFill("solid", fgColor="E2EFDA"),
            ".pptx": PatternFill("solid", fgColor="DDEBF7"),
            ".pdf": PatternFill("solid", fgColor="FFF2CC"),
            ".csv": PatternFill("solid", fgColor="FCE4D6"),
        }

        # ===== 创建全新工作簿，覆盖现有文件 =====
        wb = Workbook()
        # 移除默认的Sheet
        if "Sheet" in wb.sheetnames:
            wb.remove(wb["Sheet"])

        # ===== API 数据导出 =====
        for sheet_type in ["request", "response"]:
            sheet_name = (
                "请求详情识别结果" if sheet_type == "request" else "响应详情识别结果"
            )
            ws = wb.create_sheet(sheet_name)

            # 添加表头
            headers = [
                "编号",
                "名称",
                "类型",
                "start_line",
                "headers",
                "body",
                "匹配数量",
                "状态",
            ]
            ws.append(headers)

            # 设置表头样式
            for col in range(1, len(headers) + 1):
                cell = ws.cell(1, col)
                cell.fill = header_fill
                cell.font = header_font
                cell.alignment = Alignment(horizontal="center", vertical="center")
                cell.border = thin_border

            for item in api_asset_data_label_test_result:
                if sheet_type not in item:
                    continue
                sheet_data = item[sheet_type]
                types = ["样本", "已匹配", "未匹配", "误匹配"]
                key_map = {
                    "样本": "sample",
                    "已匹配": "matched",
                    "未匹配": "unmatched",
                    "误匹配": "misidentification",
                }
                start_row = ws.max_row + 1

                for t in types:
                    row_data = sheet_data.get(key_map[t], {})
                    row = [
                        handle_empty(sheet_data.get("id")),
                        handle_empty(sheet_data.get("name")),
                        t,
                        handle_empty(row_data.get("start_line")),
                        handle_empty(row_data.get("headers")),
                        handle_empty(row_data.get("body")),
                        handle_empty(row_data.get("count", 0)),
                        handle_empty(sheet_data.get("status")),
                    ]
                    ws.append(row)
                    cur_row = ws.max_row

                    # 设置样式
                    for col in range(1, 9):
                        cell = ws.cell(cur_row, col)
                        cell.border = thin_border
                        if col not in [3, 7, 8]:
                            cell.fill = (
                                even_row_fill if cur_row % 2 == 0 else odd_row_fill
                            )

                    ws.cell(cur_row, 3).fill = api_type_fills[t]
                    ws.cell(cur_row, 3).alignment = Alignment(
                        horizontal="center", vertical="center"
                    )

                    count_cell = ws.cell(cur_row, 7)
                    try:
                        if count_cell.value not in ("-", None, ""):
                            count_cell.font = (
                                count_blue_font
                                if int(count_cell.value) > 0
                                else count_gray_font
                            )
                        else:
                            count_cell.font = count_gray_font
                    except:
                        count_cell.font = count_gray_font
                    count_cell.alignment = Alignment(
                        horizontal="center", vertical="center"
                    )

                    status_cell = ws.cell(cur_row, 8)
                    if status_cell.value == "PASS":
                        status_cell.fill, status_cell.font = pass_fill, pass_font
                    elif status_cell.value == "FAILED":
                        status_cell.fill, status_cell.font = fail_fill, fail_font
                    status_cell.alignment = Alignment(
                        horizontal="center", vertical="center"
                    )

                # 合并编号、名称、状态
                end_row = ws.max_row
                for col in [1, 2, 8]:
                    ws.merge_cells(
                        start_row=start_row,
                        start_column=col,
                        end_row=end_row,
                        end_column=col,
                    )
                    merged_cell = ws.cell(start_row, col)
                    merged_cell.alignment = Alignment(
                        horizontal="center", vertical="center"
                    )
                    merged_cell.border = thin_border
                    if col != 8:
                        merged_cell.fill = (
                            even_row_fill if start_row % 2 == 0 else odd_row_fill
                        )
                    else:
                        if str(merged_cell.value or "").upper() == "PASS":
                            merged_cell.fill, merged_cell.font = pass_fill, pass_font
                        elif str(merged_cell.value or "").upper() == "FAILED":
                            merged_cell.fill, merged_cell.font = fail_fill, fail_font

        # ===== 文件标签数据导出 =====
        sheet_name = "文件标签识别结果"
        ws = wb.create_sheet(sheet_name)

        # 添加表头
        headers = [
            "编号",
            "名称",
            "文件类型",
            "目标文件",
            "预期数量",
            "匹配数量",
            "误匹配详情",
            "状态",
        ]
        ws.append(headers)

        # 设置表头样式
        for col in range(1, len(headers) + 1):
            cell = ws.cell(1, col)
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = Alignment(horizontal="center", vertical="center")
            cell.border = thin_border

        file_types = list(file_type_fills.keys())
        for item in file_asset_data_label_test_result:
            start_row = ws.max_row + 1
            for ft in file_types:
                file_data = item.get(ft)
                row = [
                    handle_empty(item.get("id")),
                    handle_empty(item.get("name")),
                    ft,
                    handle_empty(file_data.get("target_file") if file_data else None),
                    handle_empty(
                        file_data.get("expected_count") if file_data else None
                    ),
                    handle_empty(file_data.get("matched_count") if file_data else None),
                    handle_empty(
                        file_data.get("misidentification") if file_data else None
                    ),
                    handle_empty(item.get("status")),
                ]
                ws.append(row)
                cur_row = ws.max_row

                for col in range(1, 9):
                    cell = ws.cell(cur_row, col)
                    cell.border = thin_border
                    if col not in [3, 6, 8]:
                        cell.fill = even_row_fill if cur_row % 2 == 0 else odd_row_fill

                ws.cell(cur_row, 3).fill = file_type_fills.get(ft, odd_row_fill)
                ws.cell(cur_row, 3).alignment = Alignment(
                    horizontal="center", vertical="center"
                )

                # count 列
                count_cell = ws.cell(cur_row, 6)
                try:
                    if count_cell.value not in ("-", None, ""):
                        count_cell.font = (
                            count_blue_font
                            if int(count_cell.value) > 0
                            else count_gray_font
                        )
                    else:
                        count_cell.font = count_gray_font
                except:
                    count_cell.font = count_gray_font
                count_cell.alignment = Alignment(horizontal="right", vertical="center")

                # 状态列
                status_cell = ws.cell(cur_row, 8)
                status_value = str(status_cell.value or "").upper()
                if "PASS" in status_value:
                    status_cell.fill, status_cell.font = pass_fill, pass_font
                elif "FAILED" in status_value:
                    status_cell.fill, status_cell.font = fail_fill, fail_font
                status_cell.alignment = Alignment(
                    horizontal="center", vertical="center"
                )

            end_row = ws.max_row
            for col in [1, 2, 8]:
                ws.merge_cells(
                    start_row=start_row,
                    start_column=col,
                    end_row=end_row,
                    end_column=col,
                )
                merged_cell = ws.cell(start_row, col)
                merged_cell.alignment = Alignment(
                    horizontal="center", vertical="center"
                )
                merged_cell.border = thin_border
                if col != 8:
                    merged_cell.fill = (
                        even_row_fill if start_row % 2 == 0 else odd_row_fill
                    )
                else:
                    status_value = str(merged_cell.value or "").upper()
                    if "PASS" in status_value:
                        merged_cell.fill, merged_cell.font = pass_fill, pass_font
                    elif "FAILED" in status_value:
                        merged_cell.fill, merged_cell.font = fail_fill, fail_font

        # ===== 调整列宽和行高 =====
        sheet_columns = {
            "请求详情识别结果": {
                "A": 15,
                "B": 20,
                "C": 10,
                "D": 12,
                "E": 25,
                "F": 25,
                "G": 10,
                "H": 12,
            },
            "响应详情识别结果": {
                "A": 15,
                "B": 20,
                "C": 10,
                "D": 12,
                "E": 25,
                "F": 25,
                "G": 10,
                "H": 12,
            },
            "文件标签识别结果": {
                "A": 15,
                "B": 20,
                "C": 12,
                "D": 30,
                "E": 12,
                "F": 12,
                "G": 25,
                "H": 12,
            },
        }

        for sheet_name, widths in sheet_columns.items():
            if sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                for col, width in widths.items():
                    ws.column_dimensions[col].width = width
                for row in range(2, ws.max_row + 1):
                    ws.row_dimensions[row].height = 25

        # 保存文件（覆盖现有文件）
        wb.save(file_path)
        log.success(f"所有测试结果已导出到 {file_path}")

    async def verify_file_asset_data_label_result(
        self,
        https_req: AsyncHttpClient,
        file_asset_data_labels: List[Dict[str, Any]],
        specification_name: str,
    ):
        """验证文件资产中数据标签的识别结果

        Args:
            file_asset_data_labels (List[Dict[str, Any]]): _description_
        """
        file_asset_data_label_test_result = []

        # 由于文件详情数据存在延迟
        await ApioneUtils.is_file_asset_count_equal_expected(
            https_req, len(file_asset_data_labels) * 8
        )
        await asyncio.sleep(3)
        log.info("文件资产条目无误")

        for file_asset_data_label in file_asset_data_labels:
            file_asset_data_label_name = file_asset_data_label["name"]
            file_lists = FileUtils.get_all_files(
                folder_path=FileUtils.find_file_from_root(
                    f"files/data_label_file/test_data/{specification_name.replace('/', '_')}",
                    create_if_not_exists=True,
                ),
                file_name=f"_{file_asset_data_label_name}.",
            )
            file_asset_data_label_test_record = (
                await self.compare_file_asset_label_result(
                    https_req, file_asset_data_label, file_lists
                )
            )

            file_asset_data_label_test_result.append(file_asset_data_label_test_record)

        return file_asset_data_label_test_result

    async def verify_api_asset_data_label_result(
        self,
        https_req: AsyncHttpClient,
        app: str,
        api_asset_data_labels: List[Dict[str, Any]],
    ):
        """验证API识别的数据标签

        Args:
            api_asset_data_labels (List[Dict[str, Any]]): _description_
        """
        # 1、获取资产详情
        api_asset_data_label_test_result = []
        for api_asset_data_label in api_asset_data_labels:
            api_path = app + "/data_label_test/" + api_asset_data_label["id"]
            api_asset: Optional[ApiAssetRecord] = (
                await ApioneUtils.get_api_asset_record(
                    https_req=https_req, api=api_path
                )
            )

            api_asset_label_detail: Optional[ApiAssetLabelDetail] = (
                await ApioneUtils.get_api_asset_label_detail(
                    https_req=https_req, api_id=api_asset.id
                )
            )

            # 对比识别情况
            api_asset_data_label_test_result.append(
                self.compare_api_asset_label_result(
                    api_asset_data_label, api_asset_label_detail
                )
            )
        return api_asset_data_label_test_result

    def save_doc_file(self, record_text, doc_path):
        doc_path = Path(doc_path)
        doc_path.parent.mkdir(parents=True, exist_ok=True)

        word = win32.gencache.EnsureDispatch("Word.Application")
        word.Visible = False
        word.DisplayAlerts = 0  # wdAlertsNone

        doc = None
        try:
            doc = word.Documents.Add()
            doc.Content.Text = record_text

            # SaveAs 重试机制
            for attempt in range(3):
                try:
                    doc.SaveAs(str(doc_path), FileFormat=constants.wdFormatDocument)
                    break
                except pywintypes.com_error as e:
                    log.error(f"[Warning] SaveAs 失败 {attempt+1}/3: {e}")
                    time.sleep(1)
            else:
                log.error(f"[Error] 保存 {doc_path} 失败")

            # 关闭文档
            try:
                if doc is not None:
                    doc.Close(False)
            except pywintypes.com_error as e:
                log.error(f"[Warning] 关闭文档失败: {e}")

        finally:
            # Quit Word 时也加保护
            try:
                word.Quit()
            except pywintypes.com_error as e:
                log.error(f"[Warning] Word Quit 失败: {e}")

    def append_to_files(self, data_label_name, data_label_test_data, output_dir):
        """
        生成新文件（每次覆盖旧文件，不追加）
        生成文件：
        - CSV, DOCX, XLSX, PDF, PPTX, TXT
        - 兼容旧格式：DOC, XLS
        - 打包为ZIP
        """
        # os.makedirs(output_dir, exist_ok=True)
        base_path = Path(output_dir)

        # 公共数据
        formatted_test_data = "\n".join(data_label_test_data)
        # record_text = f"ID: {data_label_id} 名称: {data_label_name}\n测试数据:\n{formatted_test_data}\n"
        record_text = f"测试数据:\n{formatted_test_data}\n"

        files_to_zip = []

        # 1. CSV
        csv_path = base_path / f"数据标签识别_{data_label_name}.csv"
        pd.DataFrame([{"测试数据": formatted_test_data}]).to_csv(
            csv_path, index=False, encoding="utf-8"
        )
        files_to_zip.append(csv_path)

        # 2. Word (DOCX)
        docx_path = base_path / f"数据标签识别_{data_label_name}.docx"
        docx = Document()
        docx.add_paragraph(record_text)
        docx.save(docx_path)
        files_to_zip.append(docx_path)

        # 2.1 Word (DOC)
        doc_path = base_path / f"数据标签识别_{data_label_name}.doc"
        self.word_doc_manager.save_doc_file(record_text, doc_path)
        files_to_zip.append(doc_path)

        # 3. Excel (XLSX)
        xlsx_path = base_path / f"数据标签识别_{data_label_name}.xlsx"
        pd.DataFrame(
            # [{"数据标签ID": data_label_id, "数据标签名称": data_label_name, "测试数据": formatted_test_data}]
            [{"测试数据": formatted_test_data}]
        ).to_excel(xlsx_path, index=False, engine="openpyxl")
        files_to_zip.append(xlsx_path)

        # 3.1 Excel (XLS，用 csv 简单替代，兼容性写法)
        # xls_path = base_path / f"数据标签识别_{data_label_name}.xls"
        # pd.DataFrame(
        #     [{"测试数据": formatted_test_data}]
        #     # [{"数据标签ID": data_label_id, "数据标签名称": data_label_name, "测试数据": formatted_test_data}]
        # ).to_excel(xls_path, index=False, engine="xlwt")  # xlwt 支持 .xls
        # files_to_zip.append(xls_path)
        xls_path = base_path / f"数据标签识别_{data_label_name}.xls"
        workbook = xlwt.Workbook()
        sheet = workbook.add_sheet("Sheet1")

        # 写表头
        sheet.write(0, 0, "测试数据")

        # 写数据
        sheet.write(1, 0, formatted_test_data)

        # 保存文件
        workbook.save(str(xls_path))
        files_to_zip.append(xls_path)

        # 4. PDF
        pdf_path = base_path / f"数据标签识别_{data_label_name}.pdf"
        pdf = FPDF()
        pdf.add_page()
        pdf.add_font(
            "NotoSansSC",
            "",
            FileUtils.find_file_from_root(
                "testcases/test_data_label/NotoSansSC-Regular.ttf"
            ),
            uni=True,
        )
        pdf.set_font("NotoSansSC", size=12)
        pdf.multi_cell(0, 10, txt=record_text)
        pdf.output(pdf_path)
        files_to_zip.append(pdf_path)

        # 5. PPTX
        pptx_path = base_path / f"数据标签识别_{data_label_name}.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        # slide.shapes.title.text = f"数据标签识别 - {data_label_name}"
        slide.placeholders[1].text = record_text
        prs.save(pptx_path)
        files_to_zip.append(pptx_path)

        # 6. TXT
        txt_path = base_path / f"数据标签识别_{data_label_name}.txt"
        with open(txt_path, "w", encoding="utf-8") as f:
            f.write(record_text)
        files_to_zip.append(txt_path)

        # 7. 打包 ZIP
        # zip_path = base_path / f"数据标签识别_{data_label_name}.zip"
        # with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zipf:
        #     for file in files_to_zip:
        #         zipf.write(file, arcname=file.name)

        # return zip_path

    async def send_file_asset_requests(
        self,
        http_req: AsyncHttpClient,
        file_asset_data_labels: List[Dict[str, Any]],
        specification_name: str,
    ):
        """上传文件请求

        Args:
            http_req (AsyncHttpClient): _description_
            file_asset_data_labels (List[Dict[str, Any]]): _description_
        """
        # 1、构造需要上传的数据标签文件数据
        file_data_label_path = FileUtils.find_file_from_root(
            f"files/data_label_file/test_data/{specification_name.replace('/', '_')}",
            create_if_not_exists=True,
        )
        # self.generate_upload_files(file_data_label_path, file_asset_data_labels)

        # 2、发送请求
        upload_file_folder_path = "/api/upload"
        await http_req.upload_files(
            folder=file_data_label_path,
            url=upload_file_folder_path,
            use_multipart=True,
        )

    async def send_api_asset_requests(
        self, http_req: AsyncHttpClient, api_asset_data_labels: List[Dict[str, Any]]
    ):
        """发送请求

        Args:
            data_label_test_data (List[Dict[str, Any]]): _description_
        """
        url_body_maps = {
            "/data_label_test/" + data_label["id"]: data_label["body"]
            for data_label in api_asset_data_labels
        }
        log.debug(url_body_maps)
        requests = []

        for url, body in url_body_maps.items():
            for _ in range(5):  # 每个 url 发 20 次
                requests.append(
                    {
                        "method": "POST",
                        "url": url,
                        "data": json.dumps(body, ensure_ascii=False).encode("utf-8"),
                    }
                )

        # 并发执行
        await http_req.batch_request(requests, max_concurrent=5)

    def select_data_label_by_refer(
        self, data_label_refer: List[Dict[str, str]], all_data_labels: Dict[str, Dict]
    ):
        """选取指定标准下的数据标签"""
        consistent = []
        for data_label in data_label_refer:
            data_label_id = data_label["id"]
            data_label_name = data_label["name"]
            if data_label_id in all_data_labels.keys():
                # 数据标签ID相同
                base_data_label_scope = all_data_labels[data_label_id].get("scope", -1)
                base_data_label_value = all_data_labels[data_label_id].get("body", [])
                base_file_data_label_value = all_data_labels[data_label_id].get(
                    "file_data", []
                )

                consistent.append(
                    {
                        "id": data_label_id,
                        "name": data_label_name,
                        "scope": base_data_label_scope,
                        "body": base_data_label_value,
                        "file_data": base_file_data_label_value,
                    }
                )
        return consistent

    def generate_upload_files(
        self, file_data_label_path: str, file_asset_data_labels: List[Dict[str, Any]]
    ):
        # 1、查看指定文件夹是否存在，并且是否存在历史构造的数据，存在需要清理
        FileUtils.remove_all_file_in_folder(
            FileUtils.find_file_from_root(file_data_label_path)
        )

        # 2、根据file_asset_data_label 参数数据构造数据标签测试文件
        for file_asset_data_label in file_asset_data_labels:
            try:
                name = file_asset_data_label.get("name")
                file_data = file_asset_data_label.get("file_data")
                self.append_to_files(name, file_data, file_data_label_path)
            except Exception as e:
                log.exception(f"生成文件失败: {name}, 错误: {e}")
                raise

    spec_list = load_specification.__func__()

    @pytest.mark.asyncio
    @pytest.mark.parametrize(
        "specification_name,specification_id",
        [
            pytest.param(name, spec_id, id=f"{name}-id{spec_id}")
            for name, spec_id in spec_list
        ],
    )
    async def test_data_label(
        self,
        proxy_apps,
        https_req,
        http_req,
        sc_ssh_client,
        all_data_label_refers,
        all_data_labels,
        specification_name,
        specification_id,
    ):
        # 1. 获取指定标准下的标签
        await self.choose_specification(
            https_req, sc_ssh_client, specification_name, specification_id
        )
        date_label_refer = all_data_label_refers.get(specification_name)
        data_label_test_data = self.select_data_label_by_refer(
            date_label_refer, all_data_labels
        )

        # 分类
        api_asset_data_labels = [d for d in data_label_test_data if d["scope"] < 2]
        file_asset_data_labels = [
            d for d in data_label_test_data if d["scope"] % 2 == 0
        ]

        # 2. 禁用自动合并 + API测试
        await ApioneUtils.update_auto_merge_config(https_req)
        await http_req.set_url(f"http://{proxy_apps['data_label'][0]}")
        await self.send_api_asset_requests(http_req, api_asset_data_labels)
        api_asset_data_label_test_result = (
            await self.verify_api_asset_data_label_result(
                https_req, proxy_apps["data_label"][0], api_asset_data_labels
            )
        )

        # 3. 文件测试
        await http_req.set_url(f"http://{proxy_apps['data_label'][1]}")
        await self.send_file_asset_requests(
            http_req, file_asset_data_labels, specification_name
        )
        file_asset_data_label_test_result = (
            await self.verify_file_asset_data_label_result(
                https_req, file_asset_data_labels, specification_name
            )
        )

        self.export_label_to_excel(
            api_asset_data_label_test_result,
            file_asset_data_label_test_result,
            specification_name,
        )

        # --- API 统计 ---
        request_pass = sum(
            1
            for r in api_asset_data_label_test_result
            if r["request"]["status"] == "PASS"
        )
        response_pass = sum(
            1
            for r in api_asset_data_label_test_result
            if r["response"]["status"] == "PASS"
        )
        request_mis = sum(
            1
            for r in api_asset_data_label_test_result
            if r["request"]["misidentification"]["count"]
        )
        response_mis = sum(
            1
            for r in api_asset_data_label_test_result
            if r["response"]["misidentification"]["count"]
        )

        request_fail = len(api_asset_data_label_test_result) - request_pass
        response_fail = len(api_asset_data_label_test_result) - response_pass
        # --- 文件 统计 ---
        # file_types = [".doc", ".docx", ".xls", ".xlsx", ".txt", ".pptx", ".pdf", ".csv"]
        file_types = [".docx", ".xls", ".xlsx", ".txt", ".pptx", ".pdf", ".csv"]
        file_stats = {ft: {"pass": 0, "fail": 0, "mis": 0} for ft in file_types}
        total_file_pass = 0
        for record in file_asset_data_label_test_result:
            for ft in file_types:
                detail = record[ft]
                if detail["misidentification"]:
                    file_stats[ft]["mis"] += 1
                if detail["expected_count"] == detail["matched_count"]:
                    file_stats[ft]["pass"] += 1
                else:
                    file_stats[ft]["fail"] += 1
            if record["status"] == "PASS":
                total_file_pass += 1

        total_file_fail = len(file_asset_data_label_test_result) - total_file_pass

        # --- 输出结果 ---
        self.test_results.append(
            {
                "specification": specification_name,
                "total_labels": len(data_label_test_data),
                "api": {
                    "total": len(api_asset_data_labels),
                    "request_pass": request_pass,
                    "request_fail": request_fail,
                    "request_mis": request_mis,
                    "response_pass": response_pass,
                    "response_fail": response_fail,
                    "response_mis": response_mis,
                },
                "file": {
                    "total": len(file_asset_data_labels),
                    "total_pass": total_file_pass,
                    "total_fail": total_file_fail,
                    "file_stats": file_stats,
                },
            }
        )

    def export_summary_markdown(self, test_result) -> str:
        """
        导出简洁美观的 Markdown 报告
        返回格式示例：
        # 📊 数据标签测试汇总结果

        ## 🎯 标准名称
        - 总标签数: 15
        - API 标签:
        ✓ 请求位置: 通过 12 | 失败 2 | 误识别 1
        ✓ 响应位置: 通过 10 | 失败 3 | 误识别 2
        - 文件标签: 通过 8 | 失败 2
        - 文件类型明细:
            • .pdf: 通过 3 | 误识别 0
            • .docx: 通过 5 | 误识别 1
        ---
        """
        lines = []

        # 标准标题
        lines.append(f"## 🎯 {test_result['specification']} 标准数据标签测试报告")
        lines.append("")

        # 基础信息
        lines.append(f"- **总标签数**: {test_result['total_labels']}")
        lines.append("")

        # API 部分
        lines.append("- **API 标签**:")
        req_pass = test_result["api"]["request_pass"]
        req_fail = test_result["api"]["request_fail"]
        req_mis = test_result["api"]["request_mis"]

        res_pass = test_result["api"]["response_pass"]
        res_fail = test_result["api"]["response_fail"]
        res_mis = test_result["api"]["response_mis"]

        lines.append(
            f"  - 请求位置: ✓ 通过 {req_pass} | ✗ 失败 {req_fail} | ❓ 误识别 {req_mis}"
        )
        lines.append(
            f"  - 响应位置: ✓ 通过 {res_pass} | ✗ 失败 {res_fail} | ❓ 误识别 {res_mis}"
        )
        lines.append("")

        # 文件部分
        file_pass = test_result["file"]["total_pass"]
        file_fail = test_result["file"]["total_fail"]
        lines.append(f"- **文件标签**: ✓ 通过 {file_pass} | ✗ 失败 {file_fail}")

        # 文件类型明细
        if test_result["file"]["file_stats"]:
            lines.append("  - 文件类型明细:")
            for ft, stats in test_result["file"]["file_stats"].items():
                lines.append(
                    f"    • {ft}: ✓ 通过 {stats['pass']} | ✗ 失败 {stats['fail']} | ❓ 误识别 {stats['mis']}"
                )

        lines.append("")
        lines.append("---")
        lines.append("")

        return "\n".join(lines)

    @pytest.mark.asyncio
    async def test_send_notice_by_wecom_robot(self, http_req, wecom_robot):
        """发送通知

        Args:
            wecom_robot (_type_): _description_
        """
        test_data_path = FileUtils.find_file_from_root(
            "files/data_label_file/test_data"
        )
        test_result_path = FileUtils.find_file_from_root(
            "files/data_label_file/test_result"
        )
        test_data_target_path = ZipUtils.zip_files(test_data_path, "文件测试数据.zip")
        test_results_target_path = ZipUtils.zip_files(
            test_result_path, "数据标签测试结果.zip"
        )
        
        for test_result in self.test_results:
            md_report = self.export_summary_markdown(test_result)
            await wecom_robot.send_markdown(md_report)
        
        
        await http_req.set_url("http://192.192.101.156:5004")
        await http_req.upload_files(
            files=[test_results_target_path, test_data_target_path],
            url="/data_label_test",
        )

        await wecom_robot.send_text(
            content=f"测试文件&测试结果地址: http://192.192.101.156:5004/data_label_test",
            mentioned_list=["tangning", "yangquan"],
        )
